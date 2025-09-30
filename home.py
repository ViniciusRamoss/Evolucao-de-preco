import streamlit as st
import pandas as pd
import plotly.express as px
import openpyxl
from io import BytesIO
from pptx import Presentation
from pptx.util import Inches
import matplotlib.pyplot as plt
from reportlab.lib import colors
from reportlab.lib.pagesizes import letter
from reportlab.platypus import SimpleDocTemplate, Table, TableStyle, Paragraph, Spacer
from reportlab.lib.styles import getSampleStyleSheet
import tempfile
import os


# Add this helper function at the top of the file, after imports
def convert_price_string(price_str):
    if pd.isna(price_str):
        return None
    # Remove all periods except the last one (which is the decimal separator)
    if isinstance(price_str, str):
        parts = price_str.split('.')
        if len(parts) > 1:
            integer_part = ''.join(parts[:-1])
            decimal_part = parts[-1]
            price_str = f"{integer_part}.{decimal_part}"
    return float(price_str)

# Modify the load_data function
def load_data(file_path, sheet_name=None):
    # Carregar os dados do Excel
    df = pd.read_excel(file_path, sheet_name=sheet_name)
    
    # Convert price column if it exists
    if 'PRECO' in df.columns:
        df['PRECO'] = df['PRECO'].apply(convert_price_string)
    
    return df

def create_excel_report(dados_historico, dados_comparativo=None):
    output = BytesIO()
    with pd.ExcelWriter(output, engine='openpyxl') as writer:
        dados_historico.to_excel(writer, sheet_name='Histórico de Preços', index=False)
        if dados_comparativo is not None:
            dados_comparativo.to_excel(writer, sheet_name='Comparativo', index=False)
    return output.getvalue()

def create_pdf_report(dados_historico, dados_comparativo=None):
    buffer = BytesIO()
    doc = SimpleDocTemplate(buffer, pagesize=letter)
    elements = []
    styles = getSampleStyleSheet()

    # Add title
    elements.append(Paragraph("Relatório de Preços", styles['Heading1']))
    elements.append(Spacer(1, 12))

    # Add historical data table
    elements.append(Paragraph("Histórico de Preços", styles['Heading2']))
    hist_data = [dados_historico.columns.tolist()] + dados_historico.values.tolist()
    hist_table = Table(hist_data)
    hist_table.setStyle(TableStyle([
        ('BACKGROUND', (0, 0), (-1, 0), colors.grey),
        ('TEXTCOLOR', (0, 0), (-1, 0), colors.whitesmoke),
        ('ALIGN', (0, 0), (-1, -1), 'CENTER'),
        ('FONTNAME', (0, 0), (-1, 0), 'Helvetica-Bold'),
        ('FONTSIZE', (0, 0), (-1, 0), 14),
        ('BOTTOMPADDING', (0, 0), (-1, 0), 12),
        ('TEXTCOLOR', (0, 1), (-1, -1), colors.black),
        ('FONTNAME', (0, 1), (-1, -1), 'Helvetica'),
        ('FONTSIZE', (0, 1), (-1, -1), 12),
        ('GRID', (0, 0), (-1, -1), 1, colors.black)
    ]))
    elements.append(hist_table)

    if dados_comparativo is not None:
        elements.append(Spacer(1, 20))
        elements.append(Paragraph("Dados Comparativos", styles['Heading2']))
        comp_data = [dados_comparativo.columns.tolist()] + dados_comparativo.values.tolist()
        comp_table = Table(comp_data)
        comp_table.setStyle(TableStyle([
            ('BACKGROUND', (0, 0), (-1, 0), colors.grey),
            ('TEXTCOLOR', (0, 0), (-1, 0), colors.whitesmoke),
            ('ALIGN', (0, 0), (-1, -1), 'CENTER'),
            ('FONTNAME', (0, 0), (-1, 0), 'Helvetica-Bold'),
            ('FONTSIZE', (0, 0), (-1, 0), 14),
            ('BOTTOMPADDING', (0, 0), (-1, 0), 12),
            ('TEXTCOLOR', (0, 1), (-1, -1), colors.black),
            ('FONTNAME', (0, 1), (-1, -1), 'Helvetica'),
            ('FONTSIZE', (0, 1), (-1, -1), 12),
            ('GRID', (0, 0), (-1, -1), 1, colors.black)
        ]))
        elements.append(comp_table)

    doc.build(elements)
    return buffer.getvalue()

def create_ppt_report(dados_historico, fig_historico, dados_comparativo=None, fig_comparativo=None):
    try:
        prs = Presentation()
        
        # Title slide
        slide = prs.slides.add_slide(prs.slide_layouts[0])
        title = slide.shapes.title
        title.text = "Relatório de Análise de Preços"

        # Historical data slide
        slide = prs.slides.add_slide(prs.slide_layouts[5])
        title = slide.shapes.title
        title.text = "Histórico de Preços"
        
        # Save plot to temp file using a different approach
        temp_dir = tempfile.mkdtemp()
        temp_path = os.path.join(temp_dir, 'hist_plot.png')
        
        # Export figure with specific parameters
        fig_historico.write_image(
            temp_path, 
            format='png',
            engine='kaleido',
            width=1000,
            height=600,
            scale=2
        )
        
        slide.shapes.add_picture(temp_path, Inches(1), Inches(2), width=Inches(8))

        if dados_comparativo is not None and fig_comparativo is not None:
            slide = prs.slides.add_slide(prs.slide_layouts[5])
            title = slide.shapes.title
            title.text = "Análise Comparativa"
            
            comp_temp_path = os.path.join(temp_dir, 'comp_plot.png')
            fig_comparativo.write_image(
                comp_temp_path,
                format='png',
                engine='kaleido',
                width=1000,
                height=600,
                scale=2
            )
            slide.shapes.add_picture(comp_temp_path, Inches(1), Inches(2), width=Inches(8))

        # Save to BytesIO
        pptx_buffer = BytesIO()
        prs.save(pptx_buffer)
        
        # Cleanup temporary files
        import shutil
        shutil.rmtree(temp_dir)
        
        return pptx_buffer.getvalue()
        
    except Exception as e:
        st.error(f"Erro ao criar apresentação: {str(e)}")
        return None

def main():
    try:
        caminho_arquivo = 'https://github.com/ViniciusRamoss/Evolucao-de-preco/blob/main/Data.XLSM'  
        aba_dados = 'Data'
        aba_precos = 'Preco'

        # Carregar dados
        base_dados = load_data(caminho_arquivo, aba_dados)
        base_preco = load_data(caminho_arquivo, aba_precos)
        base_preco = pd.melt(base_preco, id_vars=['ID', 'STATUS','CATEGORIA'], 
                            var_name='MES', value_name='PRECO').dropna(subset=['PRECO'])
        
        dados_mesclados = pd.merge(base_dados, base_preco, on='ID', how='left')
        
        # Correção: st.logo para st.image
        st.image('images/logo-nissan.png')

        #Side bar com chatbot

        # Título e descrição
        col1, col2 = st.columns(2)
        with col1:
            st.image('images/logo-nissan2.png', width=20000 )
        with col2:
            st.header("Bem-vindo(a) ao Dashboard de Análise de Preços")
        st.text("Este é o seu portal de acesso rápido e dinâmico para a análise consolidada e histórica " \
        "dos preços de veículos. Nossa aplicação foi desenvolvida para transformar grandes volumes de " \
        "dados de planilhas em informações acionáveis, garantindo que você tenha uma visão clara da evolução " \
        "do mercado.")
        st.divider()
        
        #Filtragem de Categoria | Montadora | Modelo
        st.title("Histórico de Preço")

        categorias_disponiveis = dados_mesclados['CATEGORIA'].dropna().unique()
        selecao_categoria = st.selectbox("Categoria:", categorias_disponiveis, key="categoria_principal")

        col1, col2 = st.columns(2)
        with col1:
            montadoras_disponiveis = dados_mesclados[dados_mesclados['CATEGORIA'] == selecao_categoria]['Marca'].unique()
            montadora = st.selectbox("Montadora:", montadoras_disponiveis, key="montadora_principal")
        with col2:
            modelos_disponiveis = dados_mesclados[(dados_mesclados['CATEGORIA'] == selecao_categoria) & 
                                                 (dados_mesclados['Marca'] == montadora)]['Modelo'].unique()
            modelo = st.selectbox("Modelo:", modelos_disponiveis, key="modelo_principal")
            versoes_desativadas = st.checkbox("Incluir versões desativadas", key="checkbox_versoes")

        if not versoes_desativadas:
            dados_filtrados = dados_mesclados[(dados_mesclados['CATEGORIA'] == selecao_categoria) & 
                                            (dados_mesclados['Marca'] == montadora) & 
                                            (dados_mesclados['Modelo'] == modelo) & 
                                            (dados_mesclados['STATUS'] == 'Ativo')]
        else:
            dados_filtrados = dados_mesclados[(dados_mesclados['CATEGORIA'] == selecao_categoria) & 
                                            (dados_mesclados['Marca'] == montadora) & 
                                            (dados_mesclados['Modelo'] == modelo)]

        dados_filtrados = dados_filtrados[['ID', 'Marca', 'Modelo', 'Motor', 'Versão', 
                                         'Combustível', 'Veículo', 'MES', 'PRECO']]
        dados_filtrados = dados_filtrados.rename(columns={'MES': 'Mês', 'PRECO': 'Preço'})

        # Gráfico de linha para histórico de preços
        titulo_grafico = f"Histórico de Preço - {montadora} {modelo} ({selecao_categoria})"
        fig = px.line(dados_filtrados,
                      x='Mês',
                      y='Preço',
                      color="Versão",
                      title=titulo_grafico,
                      markers=True)
        
        fig.update_layout(
            xaxis_title="Mês",
            yaxis_title="Preço (R$)",
            title_x=0.5,
            legend_title="Versão",
            hovermode="x unified",
            legend=dict(
                orientation="h", 
                yanchor="top",   
                y=-0.2,          
                xanchor="center", 
                x=0.5,           
            )
        )

        fig.update_yaxes(tickprefix="R$ ", tickformat=",.2f")
        fig.update_traces(
            hovertemplate="<b>%{customdata[0]}</b><br>" +
                          "Mês: %{x}<br>" +
                          "Preço: R$ %{y:,.2f}<br><extra></extra>",
            customdata=dados_filtrados[['Versão']]
        )

        st.plotly_chart(fig, use_container_width=True)

        # Converter Mês para datetime
        dados_filtrados['Mês'] = pd.to_datetime(dados_filtrados['Mês'])
        
        # Obter últimos preços com valores não nulos
        ultimos_precos = (dados_filtrados[dados_filtrados['Preço'].notna()]
                         .sort_values('Mês')
                         .groupby('Versão')
                         .last()
                         .reset_index())
        
        tabela_precos = ultimos_precos[['Modelo', 'Versão', 'Combustível', 'Mês', 'Preço']]
        
        # Adicionar informações de status se incluir versões desativadas
        if versoes_desativadas:
            tabela_precos = tabela_precos.merge(
                dados_mesclados[['Versão', 'STATUS']].drop_duplicates(),
                on='Versão',
                how='left'
            )
            
            # Formatar data como dd/mm/aaaa
            tabela_precos['Data_Formatada'] = tabela_precos['Mês'].dt.strftime('%d/%m/%Y')
            
            # Adicionar informação de status na versão
            tabela_precos['Versão'] = tabela_precos.apply(
                lambda row: f"{row['Versão']} (Inativo - Último preço: {row['Data_Formatada']})" 
                if row['STATUS'] == 'Inativo' else row['Versão'],
                axis=1
            )
            
            tabela_precos = tabela_precos.drop(['STATUS', 'Data_Formatada'], axis=1)
        
        # Formatar data como dd/mm/aaaa
        tabela_precos['Mês'] = tabela_precos['Mês'].dt.strftime('%d/%m/%Y')
        
        # Formatar preço
        tabela_precos['Preço'] = pd.to_numeric(tabela_precos['Preço'], errors='coerce')
        tabela_precos['Preço'] = tabela_precos['Preço'].apply(lambda x: f"R$ {x:,.2f}" if pd.notnull(x) else "N/A")
        
        # Para a primeira tabela (Últimos Preços Registrados)
        # Converter preço para numérico antes de ordenar
        tabela_precos['Preço_Num'] = tabela_precos['Preço'].str.replace('R$ ', '').str.replace('.', '')\
            .str.replace(',', '.').astype(float)
        
        # Ordenar por preço em ordem decrescente
        tabela_precos = tabela_precos.sort_values('Preço_Num', ascending=False)
        
        # Remover coluna auxiliar usada para ordenação
        tabela_precos = tabela_precos.drop('Preço_Num', axis=1)

        st.subheader("Últimos Preços Registrados")
        st.dataframe(
            tabela_precos,
            column_config={
                "Modelo": "Modelo",
                "Versão": "Versão",
                "Combustível": "Combustível",
                "Mês": "Mês Referência",
                "Preço": "Último Preço"
            },
            hide_index=True
        )
        st.divider()

        #Comparativo de preço
        st.title("Comparativo de Preço")
        st.text("Compare o modelo selecionado acima com outros veículos")

        # Modelo de referência (do histórico)
        modelo_referencia = modelo
        montadora_referencia = montadora

        # Adicionar seletor de período
        col1, col2 = st.columns(2)
        with col1:
            data_inicial = st.date_input(
                "Data Inicial",
                min_value=pd.to_datetime(dados_mesclados['MES']).min().date(),
                max_value=pd.to_datetime(dados_mesclados['MES']).max().date(),
                value=pd.to_datetime(dados_mesclados['MES']).min().date()
            )
        with col2:
            data_final = st.date_input(
                "Data Final",
                min_value=pd.to_datetime(dados_mesclados['MES']).min().date(),
                max_value=pd.to_datetime(dados_mesclados['MES']).max().date(),
                value=pd.to_datetime(dados_mesclados['MES']).max().date()
            )

        # Multiselect para montadoras e modelos
        montadora_comparativo = st.multiselect(
            "Montadoras para comparação:", 
            montadoras_disponiveis,
            key="montadora_comp"
        )

        # Lista para armazenar todos os modelos selecionados
        modelos_selecionados = []
        
        # Para cada montadora selecionada, mostrar seus modelos disponíveis
        for mont in montadora_comparativo:
            modelos_disponiveis = dados_mesclados[
                (dados_mesclados['CATEGORIA'] == selecao_categoria) & 
                (dados_mesclados['Marca'] == mont)
            ]['Modelo'].unique()
            
            modelos_comp = st.multiselect(
                f"Modelos da {mont}:",
                modelos_disponiveis,
                key=f"modelo_comp_{mont}"
            )
            modelos_selecionados.extend([(mont, modelo) for modelo in modelos_comp])

        # Converter as datas para datetime
        data_inicial = pd.to_datetime(data_inicial)
        data_final = pd.to_datetime(data_final)

        # Filtrar dados de referência pelo período selecionado
        dados_referencia = dados_mesclados[
            (dados_mesclados['CATEGORIA'] == selecao_categoria) & 
            (dados_mesclados['Marca'] == montadora_referencia) & 
            (dados_mesclados['Modelo'] == modelo_referencia) &
            (pd.to_datetime(dados_mesclados['MES']) >= data_inicial) &
            (pd.to_datetime(dados_mesclados['MES']) <= data_final)
        ]

        # Criar gráfico para cada modelo selecionado
        for mont, modelo_comp in modelos_selecionados:
            # Dados do modelo de comparação
            dados_comparativo = dados_mesclados[
                (dados_mesclados['CATEGORIA'] == selecao_categoria) & 
                (dados_mesclados['Marca'] == mont) & 
                (dados_mesclados['Modelo'] == modelo_comp) &
                (pd.to_datetime(dados_mesclados['MES']) >= data_inicial) &
                (pd.to_datetime(dados_mesclados['MES']) <= data_final)
            ]

            # Combinar dados para o gráfico
            dados_combinados = pd.concat([
                dados_referencia[['Marca', 'Modelo', 'Versão', 'MES', 'PRECO']],
                dados_comparativo[['Marca', 'Modelo', 'Versão', 'MES', 'PRECO']]
            ])

            # Criar identificador único para cada versão
            dados_combinados['Versão_Completa'] = dados_combinados['Marca'] + ' ' + \
                                                 dados_combinados['Modelo'] + ' - ' + \
                                                 dados_combinados['Versão']

            # Gráfico de comparação
            titulo_comparativo = f"Comparativo: {montadora_referencia} {modelo_referencia} vs {mont} {modelo_comp}"
            fig_comp = px.line(
                dados_combinados,
                x='MES',
                y='PRECO',
                color='Versão_Completa',
                title=titulo_comparativo,
                markers=True
            )

            fig_comp.update_layout(
                xaxis_title="Mês",
                yaxis_title="Preço (R$)",
                title_x=0.5,
                legend_title="Versão",
                hovermode="x unified",
                legend=dict(
                    orientation="h",
                    yanchor="top",
                    y=-0.2,
                    xanchor="center",
                    x=0.5
                )
            )

            fig_comp.update_yaxes(tickprefix="R$ ", tickformat=",.2f")
            fig_comp.update_traces(
                hovertemplate="<b>%{customdata[0]}</b><br>" +
                             "Mês: %{x}<br>" +
                             "Preço: R$ %{y:,.2f}<br><extra></extra>",
                customdata=dados_combinados[['Versão_Completa']]
            )

            st.plotly_chart(fig_comp, use_container_width=True)

        # Tabela comparativa consolidada
        if modelos_selecionados:
            # Lista para armazenar dados de todos os modelos
            todos_dados = [dados_referencia]
            
            for mont, modelo_comp in modelos_selecionados:
                dados_comp = dados_mesclados[
                    (dados_mesclados['CATEGORIA'] == selecao_categoria) & 
                    (dados_mesclados['Marca'] == mont) & 
                    (dados_mesclados['Modelo'] == modelo_comp) &
                    (pd.to_datetime(dados_mesclados['MES']) >= data_inicial) &
                    (pd.to_datetime(dados_mesclados['MES']) <= data_final)
                ]
                todos_dados.append(dados_comp)

            # Consolidar todos os dados
            dados_consolidados = pd.concat(todos_dados)
            dados_consolidados['Versão_Completa'] = dados_consolidados['Marca'] + ' ' + \
                                                   dados_consolidados['Modelo'] + ' - ' + \
                                                   dados_consolidados['Versão']

            # Criar tabela de últimos preços
            dados_consolidados['MES'] = pd.to_datetime(dados_consolidados['MES'])
            ultimos_precos = dados_consolidados[dados_consolidados['PRECO'].notna()]\
                .sort_values('MES')\
                .groupby('Versão_Completa')\
                .last()\
                .reset_index()

            tabela_comp = ultimos_precos[['Versão_Completa', 'MES', 'PRECO']]
            tabela_comp = tabela_comp.rename(columns={
                'Versão_Completa': 'Versão',
                'MES': 'Mês',
                'PRECO': 'Preço'
            })

            # Formatação da tabela
            tabela_comp['Mês'] = tabela_comp['Mês'].dt.strftime('%d/%m/%Y')
            tabela_comp['Preço'] = tabela_comp['Preço'].apply(lambda x: f"R$ {x:,.2f}")

            # Ordenar por preço em ordem decrescente
            tabela_comp['Preço_Num'] = tabela_comp['Preço'].str.replace('R$ ', '').str.replace('.', '')\
                .str.replace(',', '.').astype(float)
            tabela_comp = tabela_comp.sort_values('Preço_Num', ascending=False)
            tabela_comp = tabela_comp.drop('Preço_Num', axis=1)

            st.subheader("Últimos Preços Registrados - Comparativo")
            st.dataframe(
                tabela_comp,
                hide_index=True,
                column_config={
                    "Versão": "Versão",
                    "Mês": "Mês Referência",
                    "Preço": "Último Preço"
                }
            )
        st.divider()
        #Parte de download de relatorios
        st.title("Download de Relatório")
        st.text("Gere relatórios personalizados com base nos filtros aplicados.")

        col1, col2, col3 = st.columns(3)

        with col1:
            if st.button("📥 Baixar Excel"):
                excel_data = create_excel_report(
                    dados_filtrados,
                    tabela_comp if modelos_selecionados else None
                )
                st.download_button(
                    label="📊 Clique para baixar Excel",
                    data=excel_data,
                    file_name=f"relatorio_precos_{selecao_categoria}_{modelo}.xlsx",
                    mime="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet"
                )

        with col2:
            if st.button("📥 Baixar PDF"):
                pdf_data = create_pdf_report(
                    dados_filtrados,
                    tabela_comp if modelos_selecionados else None
                )
                st.download_button(
                    label="📄 Clique para baixar PDF",
                    data=pdf_data,
                    file_name=f"relatorio_precos_{selecao_categoria}_{modelo}.pdf",
                    mime="application/pdf"
                )

        with col3:
            if st.button("📥 Baixar PowerPoint"):
                ppt_data = create_ppt_report(
                    dados_filtrados,
                    fig,
                    tabela_comp if modelos_selecionados else None,
                    fig_comp if modelos_selecionados else None
                )
                st.download_button(
                    label="📊 Clique para baixar PowerPoint",
                    data=ppt_data,
                    file_name=f"relatorio_precos_{selecao_categoria}_{modelo}.pptx",
                    mime="application/vnd.openxmlformats-officedocument.presentationml.presentation"
                )
        
    except Exception as e:
        st.error(f"Erro ao carregar os dados: {str(e)}")
        return

    # Adicionar verificação se há dados disponíveis
    if dados_mesclados.empty:
        st.warning("Não há dados disponíveis para exibição.")
        return

if __name__ == "__main__":
    main()

